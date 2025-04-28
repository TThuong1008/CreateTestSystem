# services/document_extractor.py
import PyPDF2
import logging
from typing import List, Optional
import os

class DocumentExtractor:
    def __init__(self):
        # Fallback logging if custom setup fails
        try:
            from utils.logging_config import setup_logging
            self.logger = setup_logging(__name__)
        except ImportError:
            # Use standard logging if custom setup is unavailable
            self.logger = logging.getLogger(__name__)
            self.logger.setLevel(logging.INFO)
            handler = logging.StreamHandler()
            formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
            handler.setFormatter(formatter)
            self.logger.addHandler(handler)
    
    def extract_text(self, file_path: str) -> str:
        # Get file extension
        _, file_extension = os.path.splitext(file_path)
        file_extension = file_extension.lower()
        
        try:
            # Route to appropriate extraction method based on file extension
            if file_extension == '.pdf':
                return self._extract_from_pdf(file_path)
            elif file_extension == '.docx':
                return self._extract_from_docx(file_path)
            elif file_extension == '.pptx':
                return self._extract_from_pptx(file_path)
            else:
                error_msg = f"Unsupported file format: {file_extension}"
                self.logger.error(error_msg)
                raise ValueError(error_msg)
                
        except Exception as e:
            self.logger.error(f"Error extracting text from document: {e}")
            raise
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            with open(file_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                full_text = ""
                
                # Extract text from all pages
                for page in reader.pages:
                    full_text += page.extract_text() + "\n"
                
                self.logger.info(f"Successfully extracted text from PDF: {file_path}")
                return full_text
                
        except Exception as e:
            self.logger.error(f"Error extracting text from PDF: {e}")
            raise
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX file"""
        try:
            import docx
            
            doc = docx.Document(file_path)
            full_text = []
            
            # Extract text from paragraphs
            for para in doc.paragraphs:
                full_text.append(para.text)
            
            result = "\n".join(full_text)
            self.logger.info(f"Successfully extracted text from DOCX: {file_path}")
            return result
            
        except ImportError:
            self.logger.error("Python-docx library not installed. Use 'pip install python-docx' to install it.")
            raise
        except Exception as e:
            self.logger.error(f"Error extracting text from DOCX: {e}")
            raise
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            full_text = []
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        full_text.append(shape.text)
            
            result = "\n".join(full_text)
            self.logger.info(f"Successfully extracted text from PPTX: {file_path}")
            return result
            
        except ImportError:
            self.logger.error("Python-pptx library not installed. Use 'pip install python-pptx' to install it.")
            raise
        except Exception as e:
            self.logger.error(f"Error extracting text from PPTX: {e}")
            raise
    